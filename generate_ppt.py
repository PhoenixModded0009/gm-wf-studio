import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def parse_markdown(filepath):
    """Parses the Markdown file into Frontmatter and Slides."""
    with open(filepath, 'r', encoding='utf-8') as file:
        content = file.read()

    # Split frontmatter and body
    parts = content.split('---')
    if len(parts) < 3:
        print("Error: Could not find YAML frontmatter enclosed in '---'")
        sys.exit(1)
        
    frontmatter_text = parts[1].strip()
    body_text = parts[2].strip()

    # Parse Frontmatter variables
    config = {}
    for line in frontmatter_text.split('\n'):
        if ':' in line:
            key, val = line.split(':', 1)
            config[key.strip()] = val.strip()

    # Parse Slides (Split by '# ' which denotes a new slide title)
    raw_slides = body_text.split('\n# ')
    
    # Fix the first slide missing the '#' due to the split
    if raw_slides[0].startswith('# '):
        raw_slides[0] = raw_slides[0][2:]
    
    slides = []
    for raw_slide in raw_slides:
        if not raw_slide.strip():
            continue
            
        lines = raw_slide.split('\n')
        title = lines[0].strip()
        
        slide_data = {
            'title': title,
            'layout': 'Single', # Default
            'footer': '',
            'bullets': [],
            'placeholder': '',
            'notes': ''
        }

        for line in lines[1:]:
            # Match Layout
            if line.startswith('Layout:'):
                slide_data['layout'] = line.split(':', 1)[1].strip()
            # Match Footer
            elif line.startswith('Footer:'):
                slide_data['footer'] = line.split(':', 1)[1].strip()
            # Match Notes
            elif line.startswith('> Notes:'):
                slide_data['notes'] = line.split(':', 1)[1].strip()
            # Match Placeholders
            elif line.startswith('[PLACEHOLDER:'):
                slide_data['placeholder'] = line.strip()
            # Match Bullets (handles indentation for sub-bullets)
            elif line.lstrip().startswith('- '):
                indent_level = len(line) - len(line.lstrip())
                level = 0 if indent_level == 0 else 1 # Simple 2-level nesting
                text = line.lstrip()[2:].strip()
                slide_data['bullets'].append({'level': level, 'text': text})

        slides.append(slide_data)

    return config, slides

def build_presentation(config, slides, output_filename="output.pptx"):
    """Builds the PowerPoint file using python-pptx."""
    
    # 1. Select Theme Template
    theme = config.get('Theme', 'Light').lower()
    template_file = 'dark_template.pptx' if 'dark' in theme else 'light_template.pptx'
    
    try:
        prs = Presentation(template_file)
    except Exception as e:
        raise ValueError(f"Could not load '{template_file}'. Please ensure you uploaded both 'light_template.pptx' and 'dark_template.pptx' to your GitHub repository.")

    # 2. Define Layout Mapping to Standard PowerPoint Slide Masters
    layout_map = {
        'Title': 0,
        'Single': 1,
        'Data_Heavy': 1,
        'Step_by_Step': 1,
        'Divider': 2,
        'Split': 3,
        'Comparison': 3,
        'PICO': 3,
        'Vitals_Grid': 3,
        'Algorithm': 5 
    }

    # 3. Generate Slides
    for slide_data in slides:
        layout_name = slide_data['layout']
        layout_idx = layout_map.get(layout_name, 1) # Fallback if typo occurs
        
        slide_layout = prs.slide_layouts[layout_idx]
        slide = prs.slides.add_slide(slide_layout)

        # Add Title
        if slide.shapes.title:
            slide.shapes.title.text = slide_data['title']

        # Add Bullets
        if slide_data['bullets'] and layout_idx in [1, 3]: 
            body_shape = slide.placeholders[1]
            tf = body_shape.text_frame
            tf.clear() 
            
            for i, bullet in enumerate(slide_data['bullets']):
                p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
                p.text = bullet['text']
                p.level = bullet['level']

        # Add Footer (Dynamic Citation)
        if slide_data['footer']:
            left = Inches(5.5)
            top = Inches(6.8)
            width = Inches(4)
            height = Inches(0.5)
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.text = slide_data['footer']
            p.font.size = Pt(12)
            p.font.italic = True
            p.alignment = PP_ALIGN.RIGHT

        # Add Speaker Notes
        if slide_data['notes']:
            notes_slide = slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            text_frame.text = slide_data['notes']

        # Console Alert for Image Placeholders
        if slide_data['placeholder']:
            print(f"Reminder: Drag and drop an asset into slide '{slide_data['title']}' for {slide_data['placeholder']}")

    # 4. Save the File
    prs.save(output_filename)
    print(f"\n✅ Success! Presentation saved as {output_filename}")
    print(f"Mode Used: {config.get('Mode', 'Standard')}")

if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python generate_ppt.py <your_markdown_file.md>")
        sys.exit(1)
        
    md_file = sys.argv[1]
    config, slides = parse_markdown(md_file)
    
    out_file = md_file.replace('.md', '.pptx')
    build_presentation(config, slides, out_file)

